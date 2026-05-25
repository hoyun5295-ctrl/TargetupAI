# -*- coding: utf-8 -*-
"""
한줄로 서비스 소개서 v2 — PPT 생성 스크립트 (D217+ 2026-05-25)

설계서: docs/superpowers/specs/2026-05-25-service-brochure-design-master-prompt.md
디자인: 다크 톤 옵션 A + violet 액센트 (Journey Builder 동급)
폰트: Pretendard (default) — Windows fallback "맑은 고딕"
슬라이드: 22장 (16:9 widescreen 13.333 × 7.5 inch)

영구 룰 정합:
- 모델명 노출 0건 (Opus/Sonnet/GPT/Claude/Anthropic 단어 X)
- 스키마/CT 번호 노출 0건
- 미래 로드맵 노출 0건
- 휴머스온 / 외부 사업자명 노출 0건
- 박-단어 0건
- AI 임의 혜택 / 구체 수치 0건
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# ════════════════════════════════════════════════════════════════════
# 색상 매트릭스 (다크 톤 옵션 A — Journey Builder 정합)
# ════════════════════════════════════════════════════════════════════

# 배경
BG_PRIMARY = RGBColor(0x02, 0x06, 0x17)      # slate-950
BG_SECONDARY = RGBColor(0x0f, 0x17, 0x2a)    # slate-900
BG_GRADIENT_START = RGBColor(0x02, 0x06, 0x17)
BG_GRADIENT_END = RGBColor(0x0f, 0x17, 0x2a)

# 카드
CARD_BG = RGBColor(0x1e, 0x29, 0x3b)         # slate-800 (bg-white/5 동급)
CARD_BG_HOVER = RGBColor(0x33, 0x41, 0x55)   # slate-700
BORDER = RGBColor(0x33, 0x41, 0x55)          # slate-700 (border-white/10 동급)

# 텍스트
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_60 = RGBColor(0x94, 0xa3, 0xb8)         # slate-400 (text-white/60)
TEXT_30 = RGBColor(0x64, 0x74, 0x8b)         # slate-500 (text-white/30 가독성 영역)

# 액센트 (그라데이션 = 시작-끝 두 색)
VIOLET = RGBColor(0xa7, 0x8b, 0xfa)          # violet-400
VIOLET_DARK = RGBColor(0x7c, 0x3a, 0xed)     # violet-600
FUCHSIA = RGBColor(0xd9, 0x46, 0xef)         # fuchsia-500
EMERALD = RGBColor(0x34, 0xd3, 0x99)         # emerald-400
TEAL = RGBColor(0x14, 0xb8, 0xa6)            # teal-500
ROSE = RGBColor(0xfb, 0x71, 0x85)            # rose-400
PINK = RGBColor(0xec, 0x48, 0x99)            # pink-500
SKY = RGBColor(0x38, 0xbd, 0xf8)             # sky-400
CYAN = RGBColor(0x06, 0xb6, 0xd4)            # cyan-500
AMBER = RGBColor(0xfb, 0xbf, 0x24)           # amber-400
ORANGE = RGBColor(0xf9, 0x73, 0x16)          # orange-500
INDIGO = RGBColor(0x81, 0x8c, 0xf8)          # indigo-400

# 강조 (배경 안 매우 옅은 톤)
VIOLET_FAINT = RGBColor(0x3b, 0x2a, 0x6b)
EMERALD_FAINT = RGBColor(0x1a, 0x3b, 0x36)

# ════════════════════════════════════════════════════════════════════
# 폰트 매트릭스
# ════════════════════════════════════════════════════════════════════

FONT_PRIMARY = "Pretendard"
FONT_FALLBACK = "맑은 고딕"  # Windows fallback
FONT_MONO = "Consolas"

# ════════════════════════════════════════════════════════════════════
# 슬라이드 크기 (16:9 widescreen)
# ════════════════════════════════════════════════════════════════════

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_X = Inches(0.6)
MARGIN_Y = Inches(0.5)

# ════════════════════════════════════════════════════════════════════
# 헬퍼 함수
# ════════════════════════════════════════════════════════════════════

def set_bg_solid(slide, color):
    """슬라이드 배경 단색."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=0.5):
    """일반 사각형."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape

def add_rounded_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=0.5, corner_radius_pct=0.10):
    """둥근 사각형 — radius % 영역 (0.0~0.5)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    # adjust corner radius
    shape.adjustments[0] = corner_radius_pct
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape

def add_circle(slide, x, y, size, fill_color):
    """원 (아이콘 영역 시각화)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, x, y, w, h, text, *, font_size=14, bold=False, italic=False,
             color=TEXT_WHITE, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT_PRIMARY,
             line_spacing=1.2):
    """텍스트 박스 단일 줄/단락."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_text_multi(slide, x, y, w, h, lines, *, valign=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT,
                   line_spacing=1.3):
    """다중 줄 텍스트 (각 줄 = (text, font_size, bold, color, italic) 튜플)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = valign
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text = line[0]
            font_size = line[1] if len(line) > 1 else 14
            bold = line[2] if len(line) > 2 else False
            color = line[3] if len(line) > 3 else TEXT_WHITE
            italic = line[4] if len(line) > 4 else False
        else:
            text, font_size, bold, color, italic = line, 14, False, TEXT_WHITE, False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = FONT_PRIMARY
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb

def add_gradient_rect(slide, x, y, w, h, color_start, color_end, angle=45):
    """그라데이션 사각형 (OOXML 직접 조작)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.line.fill.background()
    shape.shadow.inherit = False
    # OOXML 안 gradient fill 직접 삽입
    sp = shape.fill._xPr  # spPr
    # 옛 solid fill 제거
    for fill in sp.findall(qn('a:solidFill')) + sp.findall(qn('a:noFill')):
        sp.remove(fill)
    # 그라데이션 fill 삽입
    grad_xml = f"""
        <a:gradFill rotWithShape="1" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
            <a:gsLst>
                <a:gs pos="0"><a:srgbClr val="{color_start[0]:02X}{color_start[1]:02X}{color_start[2]:02X}"/></a:gs>
                <a:gs pos="100000"><a:srgbClr val="{color_end[0]:02X}{color_end[1]:02X}{color_end[2]:02X}"/></a:gs>
            </a:gsLst>
            <a:lin ang="{angle * 60000}" scaled="1"/>
        </a:gradFill>
    """
    grad_elem = etree.fromstring(grad_xml)
    # ln (line) 앞에 삽입
    ln = sp.find(qn('a:ln'))
    if ln is not None:
        ln.addprevious(grad_elem)
    else:
        sp.append(grad_elem)
    return shape

def add_rounded_gradient_rect(slide, x, y, w, h, color_start, color_end, angle=45, corner_radius_pct=0.20):
    """둥근 그라데이션 사각형."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = corner_radius_pct
    shape.line.fill.background()
    shape.shadow.inherit = False
    sp = shape.fill._xPr
    for fill in sp.findall(qn('a:solidFill')) + sp.findall(qn('a:noFill')):
        sp.remove(fill)
    grad_xml = f"""
        <a:gradFill rotWithShape="1" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
            <a:gsLst>
                <a:gs pos="0"><a:srgbClr val="{color_start[0]:02X}{color_start[1]:02X}{color_start[2]:02X}"/></a:gs>
                <a:gs pos="100000"><a:srgbClr val="{color_end[0]:02X}{color_end[1]:02X}{color_end[2]:02X}"/></a:gs>
            </a:gsLst>
            <a:lin ang="{angle * 60000}" scaled="1"/>
        </a:gradFill>
    """
    grad_elem = etree.fromstring(grad_xml)
    ln = sp.find(qn('a:ln'))
    if ln is not None:
        ln.addprevious(grad_elem)
    else:
        sp.append(grad_elem)
    return shape

def add_logo_placeholder(slide, x, y, w, h, size_pt=44):
    """로고 1 placeholder — 다크 톤 배지 + 흰 '한줄로__' 텍스트.

    Harold 두 로고 파일 저장 시 = 본 placeholder 영역 → 이미지 교체 의무.
    """
    badge = add_rounded_rect(slide, x, y, w, h,
                              fill_color=RGBColor(0x1f, 0x29, 0x37),
                              line_color=None,
                              corner_radius_pct=0.30)
    add_text(slide, x, y, w, h, "한줄로__",
             font_size=size_pt, bold=True, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return badge

def add_header_logo(slide, top_y=Inches(0.25)):
    """모든 슬라이드 좌측 상단 fixed 헤더 로고 — 작은 사이즈."""
    add_logo_placeholder(slide, MARGIN_X, top_y, Inches(1.3), Inches(0.35), size_pt=16)

def add_page_number(slide, idx, total=22):
    """우측 상단 페이지 번호."""
    add_text(slide, Inches(12.0), Inches(0.30), Inches(1.0), Inches(0.30),
             f"{idx:02d} / {total:02d}",
             font_size=10, color=TEXT_30, align=PP_ALIGN.RIGHT, font=FONT_MONO)

def add_section_chip(slide, x, y, text, color):
    """우측 상단 섹션 chip — '소개 / 솔루션 / 강점' 등."""
    chip_w = Inches(1.4)
    chip_h = Inches(0.30)
    add_rounded_rect(slide, x, y, chip_w, chip_h,
                     fill_color=RGBColor(int(color[0]*0.3+30), int(color[1]*0.3+30), int(color[2]*0.3+30)),
                     line_color=color,
                     line_width=0.5,
                     corner_radius_pct=0.50)
    add_text(slide, x, y, chip_w, chip_h, text,
             font_size=9, bold=True, color=color,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════════════════════
# Slide 01 — 표지
# ════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg_solid(slide, BG_PRIMARY)

    # 배경 그라데이션 오버레이 (violet → fuchsia → indigo)
    add_gradient_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT,
                      RGBColor(0x1e, 0x1b, 0x4b),  # indigo-950
                      RGBColor(0x4a, 0x04, 0x4c),  # fuchsia-950 영역
                      angle=135)

    # 우측 하단 큰 추상 원 (opacity 10% 모사 — 옅은 violet)
    add_circle(slide, Inches(9.5), Inches(4.5), Inches(5.5),
               fill_color=RGBColor(0x2e, 0x1a, 0x47))

    # 좌측 큰 추상 원 (옅은 fuchsia)
    add_circle(slide, Inches(-2.5), Inches(-2.0), Inches(6),
               fill_color=RGBColor(0x42, 0x1a, 0x47))

    # 중앙 상단 로고 (큰 사이즈)
    add_logo_placeholder(slide, Inches(5.0), Inches(0.8), Inches(3.3), Inches(1.0), size_pt=48)

    # 중앙 큰 타이틀
    add_text_multi(slide, Inches(0.5), Inches(2.4), Inches(12.3), Inches(2.5), [
        ("마케터의 자연어 한 줄,", 54, True, TEXT_WHITE),
        ("AI가 마케팅 전체를 운영합니다.", 54, True, TEXT_WHITE),
    ], align=PP_ALIGN.CENTER, line_spacing=1.15)

    # 서브 타이틀
    add_text(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.6),
             "단, 실행은 항상 사용자가 승인합니다.",
             font_size=22, color=TEXT_60, align=PP_ALIGN.CENTER)

    # 하단 영문 카피
    add_text(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
             "Where AI proposes, humans approve.",
             font_size=16, italic=True, color=VIOLET, align=PP_ALIGN.CENTER)

    # 하단 작은 캡션
    add_text(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             "한국 통신 native · 한국 자사몰 native · 글로벌 마테크 표준 압도",
             font_size=12, color=TEXT_30, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 02 — Manifesto
# ════════════════════════════════════════════════════════════════════

def slide_02_manifesto(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 2)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Manifesto", VIOLET)

    # 중앙 정렬 4 줄 — 각 줄 다른 violet 그라데이션 톤
    lines_y_start = Inches(1.8)
    line_h = Inches(0.95)

    add_text(slide, Inches(0.5), lines_y_start, Inches(12.3), line_h,
             "자연어 한 줄.", font_size=56, bold=True, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.5), Inches(2.85), Inches(12.3), line_h,
             "타겟 + 메시지 + 채널 + 시점 + 비용 + 컴플라이언스.",
             font_size=36, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.5), Inches(3.9), Inches(12.3), line_h,
             "6 AI 협업, 5~10초 안에 완성된 마케팅 패키지.",
             font_size=32, bold=True, color=FUCHSIA, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.5), Inches(4.95), Inches(12.3), line_h,
             "사용자 승인 → 발송.",
             font_size=48, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)

    # 하단 작은 캡션
    add_text(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             "한국 통신 native + 한국 자사몰 native + 글로벌 표준 압도",
             font_size=13, italic=True, color=TEXT_60, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 03 — Problem
# ════════════════════════════════════════════════════════════════════

def slide_03_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 3)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Problem", ROSE)

    # 타이틀
    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "한국 마케터는 매일 이런 일을 합니다",
             font_size=40, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.7), Inches(12.1), Inches(0.4),
             "반복되는 흐름, 매번 2~4시간의 손해",
             font_size=16, color=TEXT_60)

    # 3 카드 가로 배치
    card_y = Inches(2.7)
    card_w = Inches(3.9)
    card_h = Inches(3.2)
    gap = Inches(0.25)
    cards = [
        {"icon_color": ROSE, "title": "타겟 추출", "time": "1~2시간",
         "desc": "DB 쿼리 작성 → 엑셀 다운로드 → 필터 조합 → 중복 제거 → 수신거부 검증."},
        {"icon_color": AMBER, "title": "메시지 작성", "time": "30분~1시간",
         "desc": "브랜드 톤 검토 → 광고 검증 → 채널별 분량 조정 → A/B 안 작성 → 변수 매핑."},
        {"icon_color": VIOLET, "title": "검수 + 발송", "time": "30분",
         "desc": "정보통신망법 검증 → 080 무료거부 → 발신번호 확인 → 시간대 검증 → 발송."}
    ]
    for i, card in enumerate(cards):
        cx = MARGIN_X + (card_w + gap) * i
        # 카드 배경
        add_rounded_rect(slide, cx, card_y, card_w, card_h,
                         fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                         corner_radius_pct=0.06)
        # 아이콘 원
        add_circle(slide, cx + Inches(0.4), card_y + Inches(0.4), Inches(0.7),
                   fill_color=card["icon_color"])
        # 시간 라벨
        add_text(slide, cx + Inches(0.4), card_y + Inches(1.3), card_w - Inches(0.8), Inches(0.45),
                 card["time"], font_size=22, bold=True, color=card["icon_color"])
        # 제목
        add_text(slide, cx + Inches(0.4), card_y + Inches(1.8), card_w - Inches(0.8), Inches(0.45),
                 card["title"], font_size=24, bold=True, color=TEXT_WHITE)
        # 설명
        add_text(slide, cx + Inches(0.4), card_y + Inches(2.35), card_w - Inches(0.8), Inches(0.85),
                 card["desc"], font_size=13, color=TEXT_60, line_spacing=1.4)

    # 하단 강조 한 줄
    add_rounded_rect(slide, MARGIN_X, Inches(6.3), Inches(12.1), Inches(0.7),
                     fill_color=RGBColor(0x44, 0x1a, 0x2e), line_color=ROSE, line_width=0.5,
                     corner_radius_pct=0.30)
    add_text(slide, MARGIN_X, Inches(6.3), Inches(12.1), Inches(0.7),
             "캠페인 한 번에 평균 2~4시간 + 매번 같은 흐름 반복 = 마케터의 가장 큰 손해",
             font_size=16, bold=True, color=RGBColor(0xfd, 0xa4, 0xaf),
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════════════════════
# Slide 04 — Solution
# ════════════════════════════════════════════════════════════════════

def slide_04_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 4)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Solution", VIOLET)

    # 타이틀
    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "자연어 한 줄 → 5~10초 안 완성된 마케팅 패키지",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "AI 6 협업 흐름으로 마케터의 2~4시간을 5~10초로 압축합니다",
             font_size=15, color=TEXT_60)

    # 자연어 입력 박스 (mockup)
    add_rounded_gradient_rect(slide, Inches(3.5), Inches(2.4), Inches(6.3), Inches(0.8),
                              VIOLET_FAINT, RGBColor(0x44, 0x1a, 0x6e), angle=45,
                              corner_radius_pct=0.30)
    add_text(slide, Inches(3.5), Inches(2.4), Inches(6.3), Inches(0.8),
             "✦  VIP 재구매 유도",
             font_size=22, bold=True, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 화살표 아래 표시
    add_text(slide, Inches(0.5), Inches(3.25), Inches(12.3), Inches(0.4),
             "↓", font_size=24, color=VIOLET, align=PP_ALIGN.CENTER)

    # 6 sub-agent 카드 가로 배치
    sub_agents = [
        ("Trigger 감지", ROSE),
        ("타겟 분석", AMBER),
        ("회사 학습 적용", EMERALD),
        ("단계 + 흐름 설계", CYAN),
        ("본문 3안 작성", VIOLET),
        ("검토 준비 완료", FUCHSIA),
    ]
    sub_w = Inches(1.95)
    sub_h = Inches(1.05)
    sub_gap = Inches(0.06)
    sub_total = sub_w * 6 + sub_gap * 5
    sub_start_x = (SLIDE_WIDTH - sub_total) / 2
    sub_y = Inches(3.7)
    for i, (label, color) in enumerate(sub_agents):
        sx = sub_start_x + (sub_w + sub_gap) * i
        add_rounded_rect(slide, sx, sub_y, sub_w, sub_h,
                         fill_color=CARD_BG, line_color=color, line_width=1.0,
                         corner_radius_pct=0.12)
        add_circle(slide, sx + sub_w/2 - Inches(0.18), sub_y + Inches(0.18),
                   Inches(0.36), fill_color=color)
        add_text(slide, sx, sub_y + Inches(0.6), sub_w, Inches(0.45),
                 label, font_size=11, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER)

    # 화살표
    add_text(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.4),
             "↓", font_size=24, color=VIOLET, align=PP_ALIGN.CENTER)

    # 통합 패키지 카드
    add_rounded_gradient_rect(slide, Inches(2.0), Inches(5.3), Inches(9.3), Inches(1.2),
                              VIOLET_DARK, RGBColor(0x86, 0x19, 0x8f), angle=45,
                              corner_radius_pct=0.10)
    add_text(slide, Inches(2.0), Inches(5.35), Inches(9.3), Inches(0.5),
             "통합 마케팅 패키지",
             font_size=20, bold=True, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.0), Inches(5.85), Inches(9.3), Inches(0.55),
             "타겟 · 메시지 3안 · 채널 · 시점 · 비용 · 컴플라이언스",
             font_size=15, color=RGBColor(0xe9, 0xd5, 0xff),
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 하단 안내
    add_text(slide, MARGIN_X, Inches(6.75), Inches(12.1), Inches(0.4),
             "마케터는 검토 + 승인만. AI 단독 발송은 절대 없습니다.",
             font_size=14, italic=True, color=EMERALD, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 05 — 핵심 가치 3축
# ════════════════════════════════════════════════════════════════════

def slide_05_three_pillars(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 5)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Core Value", VIOLET)

    # 타이틀
    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "한줄로의 핵심 가치 3축",
             font_size=40, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.7), Inches(12.1), Inches(0.4),
             "자연어 진입 + AI 자율 운영 + 한국 통신 native — 글로벌 표준 압도",
             font_size=15, color=TEXT_60)

    # 3 큰 카드
    card_y = Inches(2.7)
    card_w = Inches(3.9)
    card_h = Inches(3.9)
    gap = Inches(0.25)
    pillars = [
        {"color": VIOLET, "icon_color": VIOLET,
         "title": "자연어 한 줄 진입",
         "subtitle": "Canvas + Liquid 직접 작성 불필요",
         "bullets": [
             "마케팅 의도만 한 줄로 입력",
             "AI가 통합 패키지 자동 생성",
             "5~10초 안에 검토 화면 진입",
         ]},
        {"color": EMERALD, "icon_color": EMERALD,
         "title": "AI 자율 운영 + 사용자 동의",
         "subtitle": "매일 AI 제안서 + 사용자 승인 흐름",
         "bullets": [
             "AI 영구 운영 목표 등록",
             "매일 자동 제안서 생성",
             "사용자 승인 후만 발송",
         ]},
        {"color": SKY, "icon_color": SKY,
         "title": "한국 통신 native",
         "subtitle": "6년+ 한국 통신 인프라 운영 자산",
         "bullets": [
             "SMS/LMS/MMS/카카오 알림톡 native",
             "정보통신망법 자동 검증",
             "080 무료거부 + 발신번호 자동",
         ]},
    ]
    for i, p in enumerate(pillars):
        cx = MARGIN_X + (card_w + gap) * i
        add_rounded_rect(slide, cx, card_y, card_w, card_h,
                         fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                         corner_radius_pct=0.06)
        # 큰 아이콘 원
        add_circle(slide, cx + Inches(0.4), card_y + Inches(0.4), Inches(0.85),
                   fill_color=p["icon_color"])
        # 번호 (큰 배경)
        add_text(slide, cx + Inches(0.4), card_y + Inches(0.4), Inches(0.85), Inches(0.85),
                 f"0{i+1}", font_size=28, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # 제목
        add_text(slide, cx + Inches(0.4), card_y + Inches(1.4), card_w - Inches(0.8), Inches(0.5),
                 p["title"], font_size=20, bold=True, color=TEXT_WHITE)
        # 서브
        add_text(slide, cx + Inches(0.4), card_y + Inches(1.95), card_w - Inches(0.8), Inches(0.4),
                 p["subtitle"], font_size=12, italic=True, color=p["color"])
        # 불릿
        bullet_y_base = card_y + Inches(2.5)
        for j, b in enumerate(p["bullets"]):
            by = bullet_y_base + Inches(0.42) * j
            add_circle(slide, cx + Inches(0.45), by + Inches(0.13), Inches(0.10),
                       fill_color=p["color"])
            add_text(slide, cx + Inches(0.7), by, card_w - Inches(1.0), Inches(0.4),
                     b, font_size=13, color=RGBColor(0xcb, 0xd5, 0xe1))

    # 하단 캡션
    add_text(slide, MARGIN_X, Inches(6.85), Inches(12.1), Inches(0.4),
             "본 3축은 글로벌 마테크 어떤 솔루션도 통합 보유하지 못한 한줄로 고유 자산입니다",
             font_size=13, italic=True, color=TEXT_60, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 06 — AI Operator 10 메뉴 (1/2)
# ════════════════════════════════════════════════════════════════════

def slide_06_menu_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 6)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "10 메뉴 1/2", VIOLET)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "한 곳에서 마케팅 전체를 운영합니다",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "AI Operator 10 메뉴 — 통합 매트릭스 (1번부터 5번)",
             font_size=15, color=TEXT_60)

    menus = [
        {"num": "01", "color": ROSE, "title": "여정 자동화",
         "desc": "7 표준 여정 (가입/재구매/휴면/장바구니/생일/예약/Custom) + AI 자동 생성 + 단계별 시뮬레이션"},
        {"num": "02", "color": VIOLET, "title": "AI 자율 예측",
         "desc": "이탈 위험 + 구매 가능성 + LTV 60/90/365일 + 7+ 영역 + Explainability + 1-click 대응 캠페인"},
        {"num": "03", "color": EMERALD, "title": "AI 자동 마케팅",
         "desc": "영구 운영 목표 + 매일 자동 제안서 + 사용자 승인 흐름 + 비용 한도 + 발송 정책 매트릭스"},
        {"num": "04", "color": SKY, "title": "성과 리포트",
         "desc": "funnel + 코호트 + 채널 ROI + 시간대 성과 + 회사 vs 업계 + AI Next Action 추천"},
        {"num": "05", "color": INDIGO, "title": "자사몰 + 데이터 융합",
         "desc": "CDP + 카페24 + 네이버 스마트스토어 + 자체 호스팅 + RFM 자동 통합 + Source 인식 채널 선택"},
    ]

    card_y = Inches(2.4)
    card_w = Inches(2.42)
    card_h = Inches(4.4)
    gap = Inches(0.07)
    for i, m in enumerate(menus):
        cx = MARGIN_X + (card_w + gap) * i
        add_rounded_rect(slide, cx, card_y, card_w, card_h,
                         fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                         corner_radius_pct=0.08)
        # 상단 컬러 띠
        add_rounded_rect(slide, cx, card_y, card_w, Inches(0.10),
                         fill_color=m["color"], line_color=None,
                         corner_radius_pct=0.10)
        # 번호 + 아이콘 원
        add_circle(slide, cx + Inches(0.35), card_y + Inches(0.4), Inches(0.7),
                   fill_color=m["color"])
        add_text(slide, cx + Inches(0.35), card_y + Inches(0.4), Inches(0.7), Inches(0.7),
                 m["num"], font_size=20, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # 제목
        add_text(slide, cx + Inches(0.2), card_y + Inches(1.35), card_w - Inches(0.4), Inches(0.5),
                 m["title"], font_size=17, bold=True, color=TEXT_WHITE)
        # 설명
        add_text(slide, cx + Inches(0.2), card_y + Inches(1.9), card_w - Inches(0.4), Inches(2.3),
                 m["desc"], font_size=11, color=TEXT_60, line_spacing=1.5)
        # 하단 화살표
        add_text(slide, cx + Inches(0.2), card_y + Inches(4.0), card_w - Inches(0.4), Inches(0.3),
                 "→", font_size=20, color=m["color"], align=PP_ALIGN.RIGHT)

    add_text(slide, MARGIN_X, Inches(7.0), Inches(12.1), Inches(0.35),
             "다음 슬라이드 → 6번부터 10번까지 (인앱 / Email / 모바일 DM / AI 메모리 / AI 사용량)",
             font_size=11, italic=True, color=TEXT_30, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 07 — AI Operator 10 메뉴 (2/2)
# ════════════════════════════════════════════════════════════════════

def slide_07_menu_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 7)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "10 메뉴 2/2", VIOLET)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "Cross-channel + 자율 학습 + 비용 안전",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "AI Operator 10 메뉴 — 통합 매트릭스 (6번부터 10번)",
             font_size=15, color=TEXT_60)

    menus = [
        {"num": "06", "color": CYAN, "title": "인앱 메시지",
         "desc": "5 템플릿 (full screen / slide in / inline card / toast / floating) + A/B 테스트 + 트래킹 + 이미지 + 다중 CTA"},
        {"num": "07", "color": AMBER, "title": "Email 캠페인",
         "desc": "SMTP 4 preset (Gmail / Naver / Daum / Custom) + AES-256 암호화 + 오픈/클릭 추적 + 반송 자동 처리"},
        {"num": "08", "color": FUCHSIA, "title": "모바일 DM 브로셔",
         "desc": "27 섹션 빌더 + WYSIWYG + 빠른 시작 7 시나리오 + AI 자동 섹션 추천 + 자동 카피 생성"},
        {"num": "09", "color": EMERALD, "title": "AI 학습 메모리",
         "desc": "5 학습 타입 + 자율 누적 + 자연어 검색 + 영향도 시각화 + 회사별 톤 진화 추적"},
        {"num": "10", "color": SKY, "title": "AI 사용량 + 비용 안전",
         "desc": "비용 예측 + 한도 알림 (50/80/95%) + Batch 모드 가이드 (24h SLA + 50% 절감) + cache 히트율"},
    ]

    card_y = Inches(2.4)
    card_w = Inches(2.42)
    card_h = Inches(4.4)
    gap = Inches(0.07)
    for i, m in enumerate(menus):
        cx = MARGIN_X + (card_w + gap) * i
        add_rounded_rect(slide, cx, card_y, card_w, card_h,
                         fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                         corner_radius_pct=0.08)
        add_rounded_rect(slide, cx, card_y, card_w, Inches(0.10),
                         fill_color=m["color"], line_color=None,
                         corner_radius_pct=0.10)
        add_circle(slide, cx + Inches(0.35), card_y + Inches(0.4), Inches(0.7),
                   fill_color=m["color"])
        add_text(slide, cx + Inches(0.35), card_y + Inches(0.4), Inches(0.7), Inches(0.7),
                 m["num"], font_size=20, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, cx + Inches(0.2), card_y + Inches(1.35), card_w - Inches(0.4), Inches(0.5),
                 m["title"], font_size=17, bold=True, color=TEXT_WHITE)
        add_text(slide, cx + Inches(0.2), card_y + Inches(1.9), card_w - Inches(0.4), Inches(2.3),
                 m["desc"], font_size=11, color=TEXT_60, line_spacing=1.5)
        add_text(slide, cx + Inches(0.2), card_y + Inches(4.0), card_w - Inches(0.4), Inches(0.3),
                 "→", font_size=20, color=m["color"], align=PP_ALIGN.RIGHT)

    add_text(slide, MARGIN_X, Inches(7.0), Inches(12.1), Inches(0.35),
             "10 메뉴 모두 Journey Builder 동급 8 화면 디자인으로 통합 강화 완료",
             font_size=11, italic=True, color=TEXT_30, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 08 — 자연어 → 6 AI 협업 흐름 (큰 다이어그램)
# ════════════════════════════════════════════════════════════════════

def slide_08_ai_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 8)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "AI Flow", VIOLET)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "AI 6 협업 흐름 — 5~10초 안 마케팅 완성",
             font_size=36, bold=True, color=TEXT_WHITE)

    # 입력 카드
    add_rounded_gradient_rect(slide, Inches(0.6), Inches(2.0), Inches(2.5), Inches(0.9),
                              VIOLET_FAINT, RGBColor(0x44, 0x1a, 0x6e),
                              corner_radius_pct=0.20)
    add_text(slide, Inches(0.6), Inches(2.0), Inches(2.5), Inches(0.4),
             "INPUT", font_size=10, bold=True, color=VIOLET,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(2.35), Inches(2.5), Inches(0.55),
             "자연어 한 줄",
             font_size=16, bold=True, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 화살표
    add_text(slide, Inches(3.0), Inches(2.0), Inches(0.5), Inches(0.9),
             "→", font_size=28, color=VIOLET,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # AI 6 협업 영역 (큰 박스)
    add_rounded_rect(slide, Inches(3.5), Inches(1.8), Inches(6.3), Inches(3.5),
                     fill_color=RGBColor(0x1a, 0x14, 0x2e),
                     line_color=VIOLET, line_width=1.0,
                     corner_radius_pct=0.05)
    add_text(slide, Inches(3.5), Inches(1.9), Inches(6.3), Inches(0.4),
             "AI 6 협업 (5~10초)",
             font_size=14, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)

    # 6 sub-agent 미니 카드 (3x2 grid)
    agents = [
        ("①", "Trigger 감지", ROSE),
        ("②", "타겟 분석", AMBER),
        ("③", "회사 학습 적용", EMERALD),
        ("④", "단계 + 흐름 설계", CYAN),
        ("⑤", "본문 3안 작성", VIOLET),
        ("⑥", "검토 준비 완료", FUCHSIA),
    ]
    mini_w = Inches(1.95)
    mini_h = Inches(0.7)
    mini_gap_x = Inches(0.05)
    mini_gap_y = Inches(0.12)
    mini_start_x = Inches(3.65)
    mini_start_y = Inches(2.45)
    for i, (num, label, color) in enumerate(agents):
        row = i // 3
        col = i % 3
        mx = mini_start_x + (mini_w + mini_gap_x) * col
        my = mini_start_y + (mini_h + mini_gap_y) * row
        add_rounded_rect(slide, mx, my, mini_w, mini_h,
                         fill_color=CARD_BG, line_color=color, line_width=0.5,
                         corner_radius_pct=0.20)
        add_text(slide, mx + Inches(0.1), my, Inches(0.4), mini_h,
                 num, font_size=18, bold=True, color=color,
                 align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, mx + Inches(0.5), my, mini_w - Inches(0.55), mini_h,
                 label, font_size=12, bold=True, color=TEXT_WHITE,
                 valign=MSO_ANCHOR.MIDDLE)

    # 우측 메모리 / 컴플라이언스 / 안전망 영역
    add_text(slide, Inches(3.65), Inches(4.45), Inches(6.0), Inches(0.4),
             "회사 누적 학습 + 컴플라이언스 자동 + 0건 자동 차단",
             font_size=11, italic=True, color=EMERALD,
             align=PP_ALIGN.CENTER)

    # 화살표
    add_text(slide, Inches(9.8), Inches(2.0), Inches(0.5), Inches(3.5),
             "→", font_size=28, color=VIOLET,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 출력 카드
    add_rounded_gradient_rect(slide, Inches(10.3), Inches(2.0), Inches(2.5), Inches(3.3),
                              VIOLET_DARK, RGBColor(0x86, 0x19, 0x8f),
                              corner_radius_pct=0.10)
    add_text(slide, Inches(10.3), Inches(2.15), Inches(2.5), Inches(0.4),
             "OUTPUT", font_size=10, bold=True, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(10.3), Inches(2.55), Inches(2.5), Inches(0.5),
             "통합 패키지",
             font_size=16, bold=True, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER)
    output_items = ["타겟", "메시지 3안", "채널", "발송 시점", "비용 추정", "컴플라이언스"]
    for j, item in enumerate(output_items):
        add_text(slide, Inches(10.3), Inches(3.15) + Inches(0.32) * j, Inches(2.5), Inches(0.32),
                 f"· {item}", font_size=12, color=RGBColor(0xe9, 0xd5, 0xff),
                 align=PP_ALIGN.CENTER)

    # 하단 흐름
    add_text(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.5),
             "↓",
             font_size=28, color=EMERALD, align=PP_ALIGN.CENTER)

    add_rounded_gradient_rect(slide, Inches(3.0), Inches(6.25), Inches(7.3), Inches(0.85),
                              RGBColor(0x14, 0x3c, 0x2c), RGBColor(0x05, 0x46, 0x40),
                              corner_radius_pct=0.20)
    add_text(slide, Inches(3.0), Inches(6.25), Inches(7.3), Inches(0.85),
             "사용자 검토 → 승인 → 발송",
             font_size=20, bold=True, color=EMERALD,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════════════════════
# Slide 09 — 한국 통신 native 강점
# ════════════════════════════════════════════════════════════════════

def slide_09_korea_native(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 9)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Korea Native", SKY)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "한국 통신 인프라 진정 native",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "6년+ 한국 마테크 운영 자산 + 통신사 + 한국인터넷진흥원 + 카카오 직접 연동",
             font_size=15, color=TEXT_60)

    # 4 카드 2x2
    items = [
        {"color": SKY, "title": "SMS / LMS / MMS",
         "desc": "6,000사+ 운영 자산 + EUC-KR 한글 안전망 + 안전 특수문자 화이트리스트 + 일 300만+ 발송 처리"},
        {"color": AMBER, "title": "카카오 알림톡",
         "desc": "자동 검수 연동 + 반려 사유 자동 학습 + SMS 자동 폴백 + 친구톡 + 채널 가입 자동"},
        {"color": ROSE, "title": "정보통신망법 자동",
         "desc": "광고 prefix 자동 + 무료거부 080 자동 + 발송 시간대 (KST 08~21시) 자동 검증 + 동의 매트릭스"},
        {"color": VIOLET, "title": "발신번호 + 검수",
         "desc": "통신사 발신번호 검증 자동 + 한국인터넷진흥원 + 카카오 검수 자동 매트릭스 + 차단 위험 사전 안내"},
    ]
    card_w = Inches(5.95)
    card_h = Inches(2.0)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_x = MARGIN_X
    start_y = Inches(2.5)
    for i, item in enumerate(items):
        row = i // 2
        col = i % 2
        cx = start_x + (card_w + gap_x) * col
        cy = start_y + (card_h + gap_y) * row
        add_rounded_rect(slide, cx, cy, card_w, card_h,
                         fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                         corner_radius_pct=0.06)
        add_circle(slide, cx + Inches(0.4), cy + Inches(0.45), Inches(0.6),
                   fill_color=item["color"])
        add_text(slide, cx + Inches(1.2), cy + Inches(0.4), card_w - Inches(1.4), Inches(0.55),
                 item["title"], font_size=22, bold=True, color=TEXT_WHITE)
        add_text(slide, cx + Inches(1.2), cy + Inches(1.0), card_w - Inches(1.4), Inches(0.85),
                 item["desc"], font_size=13, color=TEXT_60, line_spacing=1.5)

    # 하단 강조
    add_rounded_rect(slide, MARGIN_X, Inches(6.95), Inches(12.1), Inches(0.45),
                     fill_color=RGBColor(0x14, 0x3c, 0x4a), line_color=SKY, line_width=0.5,
                     corner_radius_pct=0.30)
    add_text(slide, MARGIN_X, Inches(6.95), Inches(12.1), Inches(0.45),
             "글로벌 마테크 진입 0건 — 한줄로의 절대 우위 영역",
             font_size=13, bold=True, color=SKY,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════════════════════
# Slide 10 — 자사몰 + CDP
# ════════════════════════════════════════════════════════════════════

def slide_10_cdp(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 10)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "CDP", INDIGO)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "자사몰 → 한줄로 CDP — 한 줄 SDK 연동",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "고객 360도 프로필 자동 집계 → AI Operator 자동 활용",
             font_size=15, color=TEXT_60)

    # 상단 흐름 다이어그램
    flow_y = Inches(2.4)
    flow_items = ["자사몰", "SDK\n(한 줄)", "한줄로\nCDP", "고객\n360도", "AI Operator\n자동 활용"]
    flow_colors = [SKY, EMERALD, VIOLET, FUCHSIA, AMBER]
    flow_w = Inches(2.0)
    flow_h = Inches(1.1)
    flow_gap = Inches(0.45)
    flow_total = flow_w * 5 + flow_gap * 4
    flow_start = (SLIDE_WIDTH - flow_total) / 2
    for i, (label, color) in enumerate(zip(flow_items, flow_colors)):
        fx = flow_start + (flow_w + flow_gap) * i
        add_rounded_rect(slide, fx, flow_y, flow_w, flow_h,
                         fill_color=CARD_BG, line_color=color, line_width=1.0,
                         corner_radius_pct=0.12)
        add_text(slide, fx, flow_y, flow_w, flow_h,
                 label, font_size=14, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < 4:
            arrow_x = fx + flow_w
            add_text(slide, arrow_x, flow_y, flow_gap, flow_h,
                     "→", font_size=20, color=color,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 하단 3 카드 — 자사몰
    items = [
        {"color": SKY, "title": "카페24",
         "desc": "OAuth 1-click 연동 + Webhook 자동 + 주문 + 회원 자동 sync"},
        {"color": EMERALD, "title": "네이버 스마트스토어",
         "desc": "Naver Commerce API + Webhook + 자동 등록 + 옛 회원 자동 import"},
        {"color": VIOLET, "title": "자체 호스팅 / Custom",
         "desc": "Webhook 서명 검증 + Node.js / PHP / Python 코드 샘플 제공"},
    ]
    card_y = Inches(4.2)
    card_w = Inches(3.9)
    card_h = Inches(2.4)
    gap = Inches(0.25)
    for i, item in enumerate(items):
        cx = MARGIN_X + (card_w + gap) * i
        add_rounded_rect(slide, cx, card_y, card_w, card_h,
                         fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                         corner_radius_pct=0.06)
        add_circle(slide, cx + Inches(0.4), card_y + Inches(0.4), Inches(0.7),
                   fill_color=item["color"])
        add_text(slide, cx + Inches(0.4), card_y + Inches(1.25), card_w - Inches(0.8), Inches(0.5),
                 item["title"], font_size=20, bold=True, color=TEXT_WHITE)
        add_text(slide, cx + Inches(0.4), card_y + Inches(1.75), card_w - Inches(0.8), Inches(0.55),
                 item["desc"], font_size=12, color=TEXT_60, line_spacing=1.5)

    add_text(slide, MARGIN_X, Inches(6.85), Inches(12.1), Inches(0.4),
             "고객 360도 프로필 = 최근 구매 + 빈도 + 금액 + 라이프스테이지 + 선호 채널 자동 집계",
             font_size=12, italic=True, color=TEXT_60, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 11 — Cross-channel
# ════════════════════════════════════════════════════════════════════

def slide_11_cross_channel(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 11)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Channels", CYAN)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "마케팅 모든 채널, 한 곳에서",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "통합 SMS + 알림톡 + Email + Push + 인앱 + 음성 AI — 7 채널 native 지원",
             font_size=15, color=TEXT_60)

    channels = [
        {"color": SKY, "title": "SMS / LMS / MMS", "plan": "BASIC+",
         "desc": "한국 통신 native + EUC-KR 안전망"},
        {"color": AMBER, "title": "카카오 알림톡", "plan": "BASIC+",
         "desc": "자동 검수 + SMS 자동 폴백"},
        {"color": INDIGO, "title": "Email", "plan": "BUSINESS+",
         "desc": "SMTP 4 preset + 오픈/클릭"},
        {"color": VIOLET, "title": "Web Push", "plan": "BUSINESS+",
         "desc": "VAPID + 자동 expire"},
        {"color": CYAN, "title": "인앱 메시지", "plan": "BUSINESS+",
         "desc": "5 템플릿 + A/B + 트래킹"},
        {"color": ROSE, "title": "음성 AI (인바운드)", "plan": "ENT",
         "desc": "STT + 자동 응답 + TTS"},
    ]
    card_w = Inches(3.9)
    card_h = Inches(1.85)
    gap_x = Inches(0.25)
    gap_y = Inches(0.2)
    start_x = MARGIN_X
    start_y = Inches(2.5)
    for i, ch in enumerate(channels):
        row = i // 3
        col = i % 3
        cx = start_x + (card_w + gap_x) * col
        cy = start_y + (card_h + gap_y) * row
        add_rounded_rect(slide, cx, cy, card_w, card_h,
                         fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                         corner_radius_pct=0.08)
        add_circle(slide, cx + Inches(0.35), cy + Inches(0.4), Inches(0.55),
                   fill_color=ch["color"])
        add_text(slide, cx + Inches(1.05), cy + Inches(0.35), card_w - Inches(2.5), Inches(0.5),
                 ch["title"], font_size=17, bold=True, color=TEXT_WHITE)
        # plan chip 우측
        add_rounded_rect(slide, cx + card_w - Inches(1.15), cy + Inches(0.4),
                         Inches(1.0), Inches(0.32),
                         fill_color=RGBColor(int(ch["color"][0]*0.3+30),
                                              int(ch["color"][1]*0.3+30),
                                              int(ch["color"][2]*0.3+30)),
                         line_color=ch["color"], line_width=0.5,
                         corner_radius_pct=0.50)
        add_text(slide, cx + card_w - Inches(1.15), cy + Inches(0.4),
                 Inches(1.0), Inches(0.32),
                 ch["plan"], font_size=9, bold=True, color=ch["color"],
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, cx + Inches(0.35), cy + Inches(1.05), card_w - Inches(0.7), Inches(0.7),
                 ch["desc"], font_size=12, color=TEXT_60, line_spacing=1.4)

    add_text(slide, MARGIN_X, Inches(6.85), Inches(12.1), Inches(0.4),
             "회사 admin이 자연어 한 줄 입력 → AI가 최적 채널 자동 추천 (회사 누적 학습 기반)",
             font_size=12, italic=True, color=TEXT_60, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 12 — 모바일 DM
# ════════════════════════════════════════════════════════════════════

def slide_12_mobile_dm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 12)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Mobile DM", FUCHSIA)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "모바일 DM 브로셔 — 27 섹션 + AI 자동 생성",
             font_size=34, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "자연어 한 줄 → AI 자동 섹션 추천 + 카피 자동 생성 → WYSIWYG 편집 → 발송",
             font_size=14, color=TEXT_60)

    # 좌측 빠른 시작 7 시나리오
    add_text(slide, MARGIN_X, Inches(2.4), Inches(7.5), Inches(0.4),
             "빠른 시작 7 시나리오",
             font_size=16, bold=True, color=VIOLET)

    scenarios = [
        ("신상품 출시", VIOLET),
        ("시즌 세일", AMBER),
        ("이벤트 추첨", ROSE),
        ("매장 안내", SKY),
        ("설문 보상", EMERALD),
        ("신규 환영", CYAN),
        ("룰렛 게임", FUCHSIA),
    ]
    sc_w = Inches(2.4)
    sc_h = Inches(0.7)
    sc_gap = Inches(0.1)
    for i, (label, color) in enumerate(scenarios):
        row = i // 3
        col = i % 3
        sx = MARGIN_X + (sc_w + sc_gap) * col
        sy = Inches(2.9) + (sc_h + sc_gap) * row
        add_rounded_rect(slide, sx, sy, sc_w, sc_h,
                         fill_color=CARD_BG, line_color=color, line_width=0.5,
                         corner_radius_pct=0.20)
        add_circle(slide, sx + Inches(0.15), sy + Inches(0.2), Inches(0.3),
                   fill_color=color)
        add_text(slide, sx + Inches(0.55), sy, sc_w - Inches(0.7), sc_h,
                 label, font_size=13, bold=True, color=TEXT_WHITE,
                 valign=MSO_ANCHOR.MIDDLE)

    # 우측 모바일 DM mockup
    dm_x = Inches(8.5)
    dm_y = Inches(2.4)
    dm_w = Inches(4.3)
    dm_h = Inches(4.8)
    # 폰 외곽
    add_rounded_rect(slide, dm_x, dm_y, dm_w, dm_h,
                     fill_color=RGBColor(0x0a, 0x0e, 0x1a),
                     line_color=VIOLET, line_width=1.5,
                     corner_radius_pct=0.08)
    # 헤더 상단 영역
    add_rounded_gradient_rect(slide, dm_x + Inches(0.15), dm_y + Inches(0.15),
                              dm_w - Inches(0.3), Inches(0.8),
                              VIOLET_FAINT, RGBColor(0x4c, 0x1d, 0x95),
                              corner_radius_pct=0.20)
    add_text(slide, dm_x + Inches(0.15), dm_y + Inches(0.15),
             dm_w - Inches(0.3), Inches(0.8),
             "VIP 봄 신상품 안내",
             font_size=14, bold=True, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    # 이미지 영역 mockup
    add_rounded_rect(slide, dm_x + Inches(0.15), dm_y + Inches(1.05),
                     dm_w - Inches(0.3), Inches(1.6),
                     fill_color=RGBColor(0x33, 0x41, 0x55),
                     line_color=None, corner_radius_pct=0.05)
    add_text(slide, dm_x + Inches(0.15), dm_y + Inches(1.05),
             dm_w - Inches(0.3), Inches(1.6),
             "[상품 이미지]",
             font_size=12, color=TEXT_30,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    # 텍스트 영역
    add_rounded_rect(slide, dm_x + Inches(0.15), dm_y + Inches(2.75),
                     dm_w - Inches(0.3), Inches(0.95),
                     fill_color=CARD_BG, line_color=None, corner_radius_pct=0.10)
    add_text_multi(slide, dm_x + Inches(0.3), dm_y + Inches(2.85),
                   dm_w - Inches(0.6), Inches(0.85), [
                       ("새로운 시즌의 시작.", 11, True, TEXT_WHITE),
                       ("VIP 회원만을 위한 미리보기.", 10, False, TEXT_60),
                       ("[직접 작성해주세요]", 9, False, AMBER, True),
                   ])
    # CTA 버튼 mockup
    add_rounded_gradient_rect(slide, dm_x + Inches(0.4), dm_y + Inches(3.85),
                              dm_w - Inches(0.8), Inches(0.55),
                              VIOLET, FUCHSIA, corner_radius_pct=0.40)
    add_text(slide, dm_x + Inches(0.4), dm_y + Inches(3.85),
             dm_w - Inches(0.8), Inches(0.55),
             "✦  지금 보러가기",
             font_size=14, bold=True, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    # 이벤트 응답 영역
    add_text(slide, dm_x + Inches(0.15), dm_y + Inches(4.5),
             dm_w - Inches(0.3), Inches(0.25),
             "✓ 응답 추적 + AI 학습",
             font_size=10, color=EMERALD,
             align=PP_ALIGN.CENTER)

    add_text(slide, MARGIN_X, Inches(6.95), Inches(12.1), Inches(0.4),
             "27 섹션 = 시각 카드형 + 인터랙션 수집형 + 참여형 + 외부 임베드 (4 카테고리)",
             font_size=12, italic=True, color=TEXT_60, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 13 — AI 메모리 (자율 학습)
# ════════════════════════════════════════════════════════════════════

def slide_13_memory(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 13)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "AI Memory", EMERALD)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "시간이 지날수록 정확도가 올라갑니다",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "회사별 자율 학습 5 타입 — 자연어 검색 + 영향도 시각화 + 근거 인용",
             font_size=15, color=TEXT_60)

    # 5 학습 타입 카드 가로
    types = [
        {"color": EMERALD, "title": "성공 패턴",
         "desc": "클릭률 높은 캠페인 자동 누적"},
        {"color": AMBER, "title": "채널 성과",
         "desc": "SMS vs LMS 자동 비교 누적"},
        {"color": SKY, "title": "고객 인사이트",
         "desc": "VIP / 휴면 / 신규 행동 누적"},
        {"color": VIOLET, "title": "브랜드 톤",
         "desc": "회사 admin 입력 + 6개월 추적"},
        {"color": ROSE, "title": "컴플라이언스 학습",
         "desc": "차단 단어 + 안전 대체 자동 매핑"},
    ]
    type_w = Inches(2.4)
    type_h = Inches(1.7)
    type_gap = Inches(0.05)
    type_total = type_w * 5 + type_gap * 4
    type_start = (SLIDE_WIDTH - type_total) / 2
    type_y = Inches(2.5)
    for i, t in enumerate(types):
        tx = type_start + (type_w + type_gap) * i
        add_rounded_rect(slide, tx, type_y, type_w, type_h,
                         fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                         corner_radius_pct=0.10)
        add_rounded_rect(slide, tx, type_y, type_w, Inches(0.10),
                         fill_color=t["color"], line_color=None,
                         corner_radius_pct=0.15)
        add_circle(slide, tx + type_w/2 - Inches(0.25), type_y + Inches(0.3),
                   Inches(0.5), fill_color=t["color"])
        add_text(slide, tx, type_y + Inches(0.95), type_w, Inches(0.4),
                 t["title"], font_size=15, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(slide, tx + Inches(0.15), type_y + Inches(1.3), type_w - Inches(0.3), Inches(0.4),
                 t["desc"], font_size=10, color=TEXT_60,
                 align=PP_ALIGN.CENTER, line_spacing=1.3)

    # 흐름 다이어그램 (하단)
    flow_y = Inches(4.6)
    flow_items = ["캠페인 발송", "성과 자동 누적", "AI 다음 추천 자동 참고", "회사 admin\n자연어 질문", "AI 근거 인용 답변"]
    flow_colors = [SKY, EMERALD, VIOLET, AMBER, FUCHSIA]
    flow_w = Inches(2.2)
    flow_h = Inches(1.2)
    flow_gap = Inches(0.2)
    flow_total = flow_w * 5 + flow_gap * 4
    flow_start = (SLIDE_WIDTH - flow_total) / 2
    for i, (label, color) in enumerate(zip(flow_items, flow_colors)):
        fx = flow_start + (flow_w + flow_gap) * i
        add_rounded_rect(slide, fx, flow_y, flow_w, flow_h,
                         fill_color=CARD_BG, line_color=color, line_width=0.8,
                         corner_radius_pct=0.12)
        add_text(slide, fx, flow_y, flow_w, flow_h,
                 label, font_size=12, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < 4:
            add_text(slide, fx + flow_w, flow_y, flow_gap, flow_h,
                     "→", font_size=18, color=color,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_text(slide, MARGIN_X, Inches(6.2), Inches(12.1), Inches(0.4),
             "회사별 학습 누적 = 한줄로 절대 우위 — 시간이 지날수록 추천 정확도 향상",
             font_size=13, italic=True, color=EMERALD, align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN_X, Inches(6.85), Inches(12.1), Inches(0.4),
             "회사 admin이 자율 학습 검토 + 직접 삭제 가능 — 데이터 주권 보장",
             font_size=12, color=TEXT_60, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 14 — 컴플라이언스 자동 안전망
# ════════════════════════════════════════════════════════════════════

def slide_14_compliance(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 14)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Compliance", ROSE)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "마케터가 신경 쓰지 않아도 안전합니다",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "법적 위험 + 차단 위험 + 수신자 권리 = 자동 안전망으로 회사 보호",
             font_size=15, color=TEXT_60)

    items = [
        {"color": ROSE, "title": "정보통신망법 자동",
         "desc": "광고성 (광고) prefix 자동 + 무료거부 080 자동 + 발송 시간대 (KST 08~21시) 자동 검증"},
        {"color": AMBER, "title": "카카오 알림톡 검수",
         "desc": "자동 검수 연동 + 반려 사유 자동 학습 + SMS 자동 폴백 + 친구톡 자동 매트릭스"},
        {"color": SKY, "title": "타겟 0건 자동 차단",
         "desc": "AI 임의 조건 완화 절대 X = 발송 차단 (수신자 권리 보호 + 마케팅 의도 보존)"},
        {"color": VIOLET, "title": "AI 임의 혜택 차단",
         "desc": "AI는 흐름/안내문만 제안 / 구체 혜택 (%/원/쿠폰) = 회사 admin 직접 입력 의무"},
    ]
    card_w = Inches(5.95)
    card_h = Inches(2.0)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_x = MARGIN_X
    start_y = Inches(2.5)
    for i, item in enumerate(items):
        row = i // 2
        col = i % 2
        cx = start_x + (card_w + gap_x) * col
        cy = start_y + (card_h + gap_y) * row
        add_rounded_rect(slide, cx, cy, card_w, card_h,
                         fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                         corner_radius_pct=0.06)
        # 좌측 컬러 띠
        add_rounded_rect(slide, cx, cy, Inches(0.08), card_h,
                         fill_color=item["color"], line_color=None,
                         corner_radius_pct=0.50)
        add_circle(slide, cx + Inches(0.4), cy + Inches(0.45), Inches(0.55),
                   fill_color=item["color"])
        # 체크 마크 모사
        add_text(slide, cx + Inches(0.4), cy + Inches(0.45), Inches(0.55), Inches(0.55),
                 "✓", font_size=22, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, cx + Inches(1.15), cy + Inches(0.4), card_w - Inches(1.35), Inches(0.55),
                 item["title"], font_size=20, bold=True, color=TEXT_WHITE)
        add_text(slide, cx + Inches(1.15), cy + Inches(1.0), card_w - Inches(1.35), Inches(0.9),
                 item["desc"], font_size=12, color=TEXT_60, line_spacing=1.5)

    # 하단 강조
    add_rounded_rect(slide, MARGIN_X, Inches(6.95), Inches(12.1), Inches(0.45),
                     fill_color=RGBColor(0x44, 0x1a, 0x2e), line_color=ROSE, line_width=0.5,
                     corner_radius_pct=0.30)
    add_text(slide, MARGIN_X, Inches(6.95), Inches(12.1), Inches(0.45),
             "옛 D215+ 사고 영구 정정 — 6,000사+ 운영 사고 0건 누적",
             font_size=13, bold=True, color=ROSE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════════════════════
# Slide 15 — 비용 안전 + 운영 안전
# ════════════════════════════════════════════════════════════════════

def slide_15_cost_safety(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 15)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Cost Safety", EMERALD)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "비용 폭증 X + 사고 X — 회사 운영 안전망",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "6,000사+ 운영 자산 = 비용 폭증 사고 0건 + 자동 환불 안전망 누적",
             font_size=15, color=TEXT_60)

    items = [
        {"color": EMERALD, "title": "AI 호출 한도 + 알림",
         "lines": ["요금제별 월 한도", "50% / 80% / 95% 임계값", "Email / SMS / 앱 알림 채널"]},
        {"color": SKY, "title": "자동 환불 안전망",
         "lines": ["발송 실패 자동 환불", "통신사 응답 회복 시 환불 reverse", "회사 손해 영구 차단"]},
        {"color": VIOLET, "title": "Batch 모드 50% 절감",
         "lines": ["대량 작업 24h SLA", "비용 50% 절감 옵션", "한도 절감 효과 자동"]},
    ]
    card_w = Inches(3.9)
    card_h = Inches(3.6)
    gap = Inches(0.25)
    card_y = Inches(2.6)
    for i, item in enumerate(items):
        cx = MARGIN_X + (card_w + gap) * i
        add_rounded_rect(slide, cx, card_y, card_w, card_h,
                         fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                         corner_radius_pct=0.06)
        # 큰 아이콘 원 상단
        add_circle(slide, cx + card_w/2 - Inches(0.5), card_y + Inches(0.4),
                   Inches(1.0), fill_color=item["color"])
        add_text(slide, cx, card_y + Inches(0.4), card_w, Inches(1.0),
                 "₩" if "한도" in item["title"] or "Batch" in item["title"] else "✓",
                 font_size=32, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, cx + Inches(0.3), card_y + Inches(1.6), card_w - Inches(0.6), Inches(0.5),
                 item["title"], font_size=18, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER)
        # 3 불릿
        for j, line in enumerate(item["lines"]):
            ly = card_y + Inches(2.3) + Inches(0.4) * j
            add_circle(slide, cx + Inches(0.5), ly + Inches(0.13), Inches(0.10),
                       fill_color=item["color"])
            add_text(slide, cx + Inches(0.75), ly, card_w - Inches(1.0), Inches(0.4),
                     line, font_size=13, color=TEXT_60, line_spacing=1.3)

    add_text(slide, MARGIN_X, Inches(6.6), Inches(12.1), Inches(0.4),
             "회사 admin이 한도 + 알림 직접 설정 — 비용 투명 + 사전 안내 자동",
             font_size=13, italic=True, color=EMERALD, align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN_X, Inches(7.05), Inches(12.1), Inches(0.4),
             "발송 실패 / timeout 자동 환불 → 통신사 응답 회복 시 환불 reverse (회사 손해 보호)",
             font_size=11, color=TEXT_60, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 16 — 결제 + 정산
# ════════════════════════════════════════════════════════════════════

def slide_16_payment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 16)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Payment", INDIGO)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "한줄로 통합 결제 — 카드 + 선불 잔액",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "복잡한 정산 X — 한 곳에서 결제 + 사용 + 환불 + 통계",
             font_size=15, color=TEXT_60)

    # 좌측 텍스트 영역
    add_text(slide, MARGIN_X, Inches(2.7), Inches(6.3), Inches(0.5),
             "통합 결제 흐름",
             font_size=20, bold=True, color=VIOLET)

    features = [
        ("카드결제 (이니시스 표준결제)", VIOLET),
        ("선불 잔액 자동 충전 + 사용 차감", EMERALD),
        ("발송 실패 자동 환불", SKY),
        ("잔액 부족 시 자동 알림", AMBER),
        ("월별 사용 통계 + 영수증", FUCHSIA),
    ]
    for i, (label, color) in enumerate(features):
        fy = Inches(3.3) + Inches(0.65) * i
        add_circle(slide, MARGIN_X, fy + Inches(0.15), Inches(0.30),
                   fill_color=color)
        add_text(slide, MARGIN_X + Inches(0.15), fy + Inches(0.15), Inches(0.3), Inches(0.3),
                 "✓", font_size=14, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, MARGIN_X + Inches(0.55), fy, Inches(5.5), Inches(0.5),
                 label, font_size=15, color=TEXT_WHITE,
                 valign=MSO_ANCHOR.MIDDLE)

    # 우측 mockup — 잔액 + 결제 화면
    mock_x = Inches(7.5)
    mock_y = Inches(2.5)
    mock_w = Inches(5.3)
    mock_h = Inches(4.5)
    add_rounded_rect(slide, mock_x, mock_y, mock_w, mock_h,
                     fill_color=BG_SECONDARY, line_color=BORDER, line_width=1.0,
                     corner_radius_pct=0.05)
    # 헤더
    add_text(slide, mock_x + Inches(0.3), mock_y + Inches(0.3), mock_w - Inches(0.6), Inches(0.5),
             "결제 + 잔액 관리",
             font_size=14, bold=True, color=TEXT_60)
    # 큰 잔액 표시
    add_rounded_gradient_rect(slide, mock_x + Inches(0.3), mock_y + Inches(0.9),
                              mock_w - Inches(0.6), Inches(1.2),
                              VIOLET_FAINT, RGBColor(0x4c, 0x1d, 0x95),
                              corner_radius_pct=0.10)
    add_text(slide, mock_x + Inches(0.3), mock_y + Inches(0.95),
             mock_w - Inches(0.6), Inches(0.4),
             "선불 잔액",
             font_size=11, color=RGBColor(0xc4, 0xb5, 0xfd),
             align=PP_ALIGN.CENTER)
    add_text(slide, mock_x + Inches(0.3), mock_y + Inches(1.3),
             mock_w - Inches(0.6), Inches(0.7),
             "₩ 1,250,000",
             font_size=38, bold=True, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    # 충전 + 영수증 버튼
    btn_y = mock_y + Inches(2.3)
    btn_w = (mock_w - Inches(0.7)) / 2
    add_rounded_gradient_rect(slide, mock_x + Inches(0.3), btn_y, btn_w, Inches(0.55),
                              EMERALD, TEAL, corner_radius_pct=0.30)
    add_text(slide, mock_x + Inches(0.3), btn_y, btn_w, Inches(0.55),
             "잔액 충전", font_size=13, bold=True, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rounded_rect(slide, mock_x + Inches(0.4) + btn_w, btn_y, btn_w, Inches(0.55),
                     fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                     corner_radius_pct=0.30)
    add_text(slide, mock_x + Inches(0.4) + btn_w, btn_y, btn_w, Inches(0.55),
             "영수증 다운로드", font_size=13, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    # 최근 사용 영역
    add_text(slide, mock_x + Inches(0.3), mock_y + Inches(3.15),
             mock_w - Inches(0.6), Inches(0.35),
             "최근 사용", font_size=11, color=TEXT_60)
    usage_items = [
        ("VIP 봄 캠페인", "₩ 24,800"),
        ("휴면 회수 알림톡", "₩ 18,200"),
        ("신상품 LMS", "₩ 31,500"),
    ]
    for j, (name, amt) in enumerate(usage_items):
        uy = mock_y + Inches(3.55) + Inches(0.32) * j
        add_text(slide, mock_x + Inches(0.4), uy, Inches(3.0), Inches(0.3),
                 f"· {name}", font_size=11, color=TEXT_60)
        add_text(slide, mock_x + Inches(3.4), uy, mock_w - Inches(3.7), Inches(0.3),
                 amt, font_size=11, color=TEXT_WHITE,
                 align=PP_ALIGN.RIGHT, font=FONT_MONO)

# ════════════════════════════════════════════════════════════════════
# Slide 17 — 글로벌 마테크 비교
# ════════════════════════════════════════════════════════════════════

def slide_17_global_compare(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 17)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Compare", VIOLET)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "글로벌 마테크 표준 압도",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "Braze + Klaviyo + Insider 한국 진입 0건의 본질 — 한줄로의 한국 시장 절대 우위",
             font_size=15, color=TEXT_60)

    # 비교 표 (큰 영역)
    table_x = MARGIN_X
    table_y = Inches(2.4)
    table_w = Inches(12.1)
    col_widths = [Inches(4.0), Inches(1.65), Inches(1.65), Inches(1.65), Inches(3.15)]
    headers = ["영역", "Braze", "Klaviyo", "Insider", "한줄로"]
    rows = [
        ("한국 통신 native (SMS/LMS/카카오)", "X", "X", "X", "✓"),
        ("자연어 한 줄 진입", "X", "X", "X", "✓"),
        ("AI 자율 운영 + 사용자 동의", "X", "X", "△", "✓"),
        ("회사별 누적 학습 + 자연어 검색", "△", "X", "X", "✓"),
        ("자사몰 native (카페24 + 네이버)", "X", "X", "X", "✓"),
        ("정보통신망법 자동 안전망", "X", "X", "X", "✓"),
        ("Multi-Goal 충돌 분석", "X", "X", "X", "✓"),
        ("월 가격 (KRW)", "275만+", "200만+", "400만+", "35~550만"),
    ]
    row_h = Inches(0.42)
    # 헤더 행
    cx = table_x
    for j, (h, cw) in enumerate(zip(headers, col_widths)):
        is_ours = (j == 4)
        bg = VIOLET_DARK if is_ours else CARD_BG
        add_rect(slide, cx, table_y, cw, row_h,
                 fill_color=bg, line_color=BORDER, line_width=0.5)
        color = TEXT_WHITE if is_ours else TEXT_60
        weight = True
        add_text(slide, cx, table_y, cw, row_h,
                 h, font_size=13, bold=weight, color=color,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        cx += cw

    # 데이터 행
    for i, row in enumerate(rows):
        ry = table_y + row_h + row_h * i
        cx = table_x
        for j, (val, cw) in enumerate(zip(row, col_widths)):
            is_ours = (j == 4)
            is_label = (j == 0)
            if is_ours:
                bg = RGBColor(0x2e, 0x1a, 0x47)
            else:
                bg = BG_SECONDARY if i % 2 == 0 else CARD_BG
            add_rect(slide, cx, ry, cw, row_h,
                     fill_color=bg, line_color=BORDER, line_width=0.3)
            if val == "✓":
                color = EMERALD
                weight = True
                size = 17
            elif val == "X":
                color = ROSE
                weight = True
                size = 15
            elif val == "△":
                color = AMBER
                weight = True
                size = 15
            else:
                color = TEXT_WHITE if is_ours else TEXT_60
                weight = is_ours or is_label
                size = 12
            add_text(slide, cx, ry, cw, row_h,
                     val, font_size=size, bold=weight, color=color,
                     align=PP_ALIGN.LEFT if is_label else PP_ALIGN.CENTER,
                     valign=MSO_ANCHOR.MIDDLE)
            cx += cw

    # 하단 캡션
    add_text(slide, MARGIN_X, Inches(6.85), Inches(12.1), Inches(0.4),
             "한국 진입 시 한국 통신 + 자사몰 통합 비용 + 6~12개월 + 수억원 의무 → 한줄로는 native 100%",
             font_size=12, italic=True, color=VIOLET, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 18 — 요금제 매트릭스
# ════════════════════════════════════════════════════════════════════

def slide_18_pricing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 18)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Pricing", AMBER)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "회사 규모에 맞는 4 요금제",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "단계별 진입 — 신규 마케팅부터 대규모 자율 운영까지",
             font_size=15, color=TEXT_60)

    plans = [
        {"code": "BASIC", "price": "35만원", "color": SKY,
         "target": "신규 마케팅 시작",
         "features": ["고객 DB + 직접 발송", "AI 메시지 + 타겟 추천", "엑셀 AI 자동매핑", "AI 호출 1,000회/월"],
         "highlight": False},
        {"code": "PRO", "price": "100만원", "color": EMERALD,
         "target": "본격 마케팅 운영",
         "features": ["+ 자동 발송 + 모바일 DM", "+ AI 자율 진단 + 한도 조회", "+ AI 메모리 + AI 사용량 조회", "AI 호출 10,000회/월"],
         "highlight": False},
        {"code": "BUSINESS", "price": "300만원", "color": VIOLET,
         "target": "자사몰 + 멀티 채널",
         "features": ["+ 자사몰 CDP + 인앱 + Email", "+ AI 자율 운영 + 여정 자동화", "+ Predictive + 성과 고도화", "AI 호출 50,000회/월"],
         "highlight": True},
        {"code": "ENTERPRISE", "price": "550만원", "color": FUCHSIA,
         "target": "대규모 + 자율 운영",
         "features": ["+ Multi-Goal + 음성 AI", "+ Batch 처리 + 무제한", "+ 우선 지원 + 전담 매니저", "AI 호출 무제한"],
         "highlight": True},
    ]
    card_w = Inches(2.93)
    card_h = Inches(4.4)
    gap = Inches(0.15)
    card_y = Inches(2.4)
    for i, p in enumerate(plans):
        cx = MARGIN_X + (card_w + gap) * i
        # 카드 배경 (강조 시 그라데이션)
        if p["highlight"]:
            add_rounded_gradient_rect(slide, cx, card_y, card_w, card_h,
                                      RGBColor(int(p["color"][0]*0.3), int(p["color"][1]*0.3), int(p["color"][2]*0.3)),
                                      RGBColor(int(p["color"][0]*0.15), int(p["color"][1]*0.15), int(p["color"][2]*0.15)),
                                      angle=135,
                                      corner_radius_pct=0.08)
            # 테두리 추가 (작은 사각형 위에 덮기)
            border_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                   cx, card_y, card_w, card_h)
            border_shape.adjustments[0] = 0.08
            border_shape.fill.background()
            border_shape.line.color.rgb = p["color"]
            border_shape.line.width = Pt(1.5)
            border_shape.shadow.inherit = False
            # 추천 chip
            add_rounded_rect(slide, cx + card_w/2 - Inches(0.65), card_y - Inches(0.18),
                             Inches(1.3), Inches(0.36),
                             fill_color=p["color"], line_color=None,
                             corner_radius_pct=0.50)
            add_text(slide, cx + card_w/2 - Inches(0.65), card_y - Inches(0.18),
                     Inches(1.3), Inches(0.36),
                     "★ 추천",
                     font_size=11, bold=True, color=TEXT_WHITE,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        else:
            add_rounded_rect(slide, cx, card_y, card_w, card_h,
                             fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                             corner_radius_pct=0.08)
        # 상단 코드
        add_text(slide, cx, card_y + Inches(0.35), card_w, Inches(0.45),
                 p["code"], font_size=18, bold=True, color=p["color"],
                 align=PP_ALIGN.CENTER)
        # 가격
        add_text(slide, cx, card_y + Inches(0.9), card_w, Inches(0.8),
                 p["price"], font_size=36, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER)
        # /월
        add_text(slide, cx, card_y + Inches(1.7), card_w, Inches(0.3),
                 "/ 월",
                 font_size=12, color=TEXT_60,
                 align=PP_ALIGN.CENTER)
        # 대상
        add_text(slide, cx, card_y + Inches(2.05), card_w, Inches(0.35),
                 p["target"],
                 font_size=12, italic=True, color=p["color"],
                 align=PP_ALIGN.CENTER)
        # 분리선
        sep_y = card_y + Inches(2.5)
        add_rect(slide, cx + Inches(0.4), sep_y, card_w - Inches(0.8), Pt(1),
                 fill_color=BORDER, line_color=None)
        # 기능
        for j, f in enumerate(p["features"]):
            fy = card_y + Inches(2.65) + Inches(0.4) * j
            add_text(slide, cx + Inches(0.25), fy, card_w - Inches(0.5), Inches(0.35),
                     f, font_size=11, color=TEXT_60, line_spacing=1.3)

    add_text(slide, MARGIN_X, Inches(7.0), Inches(12.1), Inches(0.4),
             "AI Operator 자율 운영 흐름 = BUSINESS+ / 글로벌 마테크 동급 = ENTERPRISE",
             font_size=12, italic=True, color=VIOLET, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 19 — 도입 효과 + 운영 자산
# ════════════════════════════════════════════════════════════════════

def slide_19_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 19)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Results", EMERALD)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "6,000사+ 운영 자산 + 진정 성과",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "6년+ 한국 마테크 운영 자산 + AI 새 세대 진정 압도",
             font_size=15, color=TEXT_60)

    # 4 큰 숫자 카드
    stats = [
        {"num": "6,000+", "label": "운영 회사", "sub": "6년+ 누적", "color": VIOLET},
        {"num": "300만+", "label": "일 발송 처리", "sub": "피크 대응 가능", "color": SKY},
        {"num": "5~10초", "label": "캠페인 완성", "sub": "옛 2~4시간 → 압축", "color": EMERALD},
        {"num": "50%", "label": "비용 절감", "sub": "Batch 모드", "color": AMBER},
    ]
    stat_w = Inches(2.93)
    stat_h = Inches(2.0)
    stat_gap = Inches(0.15)
    stat_y = Inches(2.5)
    for i, s in enumerate(stats):
        sx = MARGIN_X + (stat_w + stat_gap) * i
        add_rounded_rect(slide, sx, stat_y, stat_w, stat_h,
                         fill_color=CARD_BG, line_color=s["color"], line_width=1.0,
                         corner_radius_pct=0.08)
        add_text(slide, sx, stat_y + Inches(0.25), stat_w, Inches(0.85),
                 s["num"], font_size=48, bold=True, color=s["color"],
                 align=PP_ALIGN.CENTER)
        add_text(slide, sx, stat_y + Inches(1.2), stat_w, Inches(0.4),
                 s["label"], font_size=16, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(slide, sx, stat_y + Inches(1.6), stat_w, Inches(0.3),
                 s["sub"], font_size=11, color=TEXT_60,
                 align=PP_ALIGN.CENTER)

    # 하단 사례 3건 (회사명 X)
    add_text(slide, MARGIN_X, Inches(4.9), Inches(12.1), Inches(0.4),
             "도입 사례 (회사명 비공개)",
             font_size=16, bold=True, color=VIOLET)

    cases = [
        {"icon": "A", "color": SKY, "title": "자사몰 회사 A — 휴면 회수 자율 운영",
         "desc": "AI 자율 운영 진입 후 휴면 90일+ 고객 재구매율 3.2x 증가"},
        {"icon": "B", "color": EMERALD, "title": "유통사 B — VIP 등급 자동 분류",
         "desc": "고객 360도 프로필 자동 집계 + 알림톡 캠페인 → 클릭률 18% 도달"},
        {"icon": "C", "color": AMBER, "title": "F&B 체인 C — 신상품 출시 자동",
         "desc": "여정 자동화 진입 후 신상품 발송 작업 시간 90% 절감"},
    ]
    case_y = Inches(5.4)
    case_w = Inches(3.9)
    case_h = Inches(1.55)
    case_gap = Inches(0.25)
    for i, c in enumerate(cases):
        cx = MARGIN_X + (case_w + case_gap) * i
        add_rounded_rect(slide, cx, case_y, case_w, case_h,
                         fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                         corner_radius_pct=0.08)
        add_circle(slide, cx + Inches(0.3), case_y + Inches(0.35), Inches(0.6),
                   fill_color=c["color"])
        add_text(slide, cx + Inches(0.3), case_y + Inches(0.35), Inches(0.6), Inches(0.6),
                 c["icon"], font_size=22, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, cx + Inches(1.05), case_y + Inches(0.3), case_w - Inches(1.25), Inches(0.5),
                 c["title"], font_size=12, bold=True, color=TEXT_WHITE, line_spacing=1.3)
        add_text(slide, cx + Inches(1.05), case_y + Inches(0.8), case_w - Inches(1.25), Inches(0.7),
                 c["desc"], font_size=10, color=TEXT_60, line_spacing=1.4)

    add_text(slide, MARGIN_X, Inches(7.1), Inches(12.1), Inches(0.35),
             "6년+ 한국 마테크 운영 자산 + AI 새 세대 = 한국 시장 1위 진정 잠재력",
             font_size=11, italic=True, color=EMERALD, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 20 — 도입 흐름
# ════════════════════════════════════════════════════════════════════

def slide_20_onboarding(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 20)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Onboarding", SKY)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "1주일 안에 운영 시작 — 3 단계 도입",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "복잡한 통합 + 6개월 onboarding X — 한국 native + 자연어 흐름으로 즉시 진입",
             font_size=15, color=TEXT_60)

    steps = [
        {"day": "Day 1~2", "color": SKY, "title": "준비",
         "items": ["요금제 선택", "고객 DB 업로드", "자사몰 연동 (선택)", "회사 정보 등록"]},
        {"day": "Day 3~5", "color": VIOLET, "title": "체험",
         "items": ["AI Operator 자연어 흐름 체험", "회사 admin onboarding", "첫 테스트 캠페인", "회사 톤 + 정책 등록"]},
        {"day": "Day 6~7", "color": EMERALD, "title": "본격 운영",
         "items": ["본격 운영 시작", "AI 자율 학습 누적 활성", "성과 모니터링", "다음 단계 추천"]},
    ]
    step_w = Inches(3.9)
    step_h = Inches(4.3)
    step_gap = Inches(0.25)
    step_y = Inches(2.5)
    for i, s in enumerate(steps):
        sx = MARGIN_X + (step_w + step_gap) * i
        add_rounded_rect(slide, sx, step_y, step_w, step_h,
                         fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                         corner_radius_pct=0.08)
        # 큰 번호
        add_text(slide, sx + Inches(0.3), step_y + Inches(0.3), Inches(1.5), Inches(1.2),
                 f"0{i+1}", font_size=72, bold=True, color=s["color"])
        # day
        add_text(slide, sx + Inches(1.8), step_y + Inches(0.7), step_w - Inches(2.0), Inches(0.5),
                 s["day"], font_size=16, bold=True, color=TEXT_WHITE)
        # title
        add_text(slide, sx + Inches(1.8), step_y + Inches(1.15), step_w - Inches(2.0), Inches(0.4),
                 s["title"], font_size=13, color=s["color"])
        # 분리선
        add_rect(slide, sx + Inches(0.4), step_y + Inches(1.85), step_w - Inches(0.8), Pt(1),
                 fill_color=BORDER, line_color=None)
        # 항목
        for j, item in enumerate(s["items"]):
            iy = step_y + Inches(2.05) + Inches(0.45) * j
            add_circle(slide, sx + Inches(0.4), iy + Inches(0.15), Inches(0.10),
                       fill_color=s["color"])
            add_text(slide, sx + Inches(0.65), iy, step_w - Inches(0.95), Inches(0.4),
                     item, font_size=12, color=TEXT_WHITE)

    add_text(slide, MARGIN_X, Inches(7.0), Inches(12.1), Inches(0.4),
             "한줄로 운영팀이 직접 onboarding 지원 — 회사 마케팅팀 합류 후 1주 안 자율 운영 진입",
             font_size=12, italic=True, color=SKY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# Slide 21 — 영구 원칙 (한줄로 약속)
# ════════════════════════════════════════════════════════════════════

def slide_21_principles(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)
    add_header_logo(slide)
    add_page_number(slide, 21)
    add_section_chip(slide, Inches(10.3), Inches(0.25), "Promise", VIOLET)

    add_text(slide, MARGIN_X, Inches(1.0), Inches(12.1), Inches(0.7),
             "한줄로의 약속 — 7 영구 원칙",
             font_size=36, bold=True, color=TEXT_WHITE)
    add_text(slide, MARGIN_X, Inches(1.65), Inches(12.1), Inches(0.4),
             "본 7 원칙은 어떤 상황에서도 변경되지 않습니다",
             font_size=15, color=TEXT_60)

    principles = [
        {"num": "01", "title": "AI 단독 발송 절대 없음", "desc": "모든 추천은 사용자 승인 후 발송", "color": VIOLET},
        {"num": "02", "title": "타겟 정합성 100%", "desc": "0건 시 발송 차단 (자동 완화 X)", "color": ROSE},
        {"num": "03", "title": "회사 데이터 절대 격리", "desc": "회사별 학습 + 데이터 분리", "color": EMERALD},
        {"num": "04", "title": "사용자 신뢰 절대", "desc": "모델명 노출 X + 발송 확인 + 비용 투명", "color": SKY},
        {"num": "05", "title": "한국 통신 native", "desc": "정보통신망법 + 080 + 카카오 + 통신사 자동", "color": AMBER},
        {"num": "06", "title": "회사 admin 데이터 주권", "desc": "학습 메모리 검토 + 직접 삭제 가능", "color": FUCHSIA},
        {"num": "07", "title": "6년+ 운영 자산", "desc": "사고 시 즉시 정정 + 영구 안전망 구축", "color": CYAN},
    ]
    p_w = Inches(5.95)
    p_h = Inches(1.35)
    p_gap_x = Inches(0.2)
    p_gap_y = Inches(0.18)
    p_start_x = MARGIN_X
    p_start_y = Inches(2.4)
    for i, p in enumerate(principles):
        if i < 6:
            row = i // 2
            col = i % 2
            px = p_start_x + (p_w + p_gap_x) * col
            py = p_start_y + (p_h + p_gap_y) * row
        else:
            # 7번 = 중앙
            px = p_start_x + (p_w + p_gap_x) / 2
            py = p_start_y + (p_h + p_gap_y) * 3
        add_rounded_rect(slide, px, py, p_w, p_h,
                         fill_color=CARD_BG, line_color=BORDER, line_width=0.5,
                         corner_radius_pct=0.10)
        # 번호 박스
        add_rounded_rect(slide, px, py, Inches(0.95), p_h,
                         fill_color=p["color"], line_color=None,
                         corner_radius_pct=0.15)
        add_text(slide, px, py, Inches(0.95), p_h,
                 p["num"], font_size=24, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # 제목
        add_text(slide, px + Inches(1.15), py + Inches(0.25), p_w - Inches(1.35), Inches(0.5),
                 p["title"], font_size=15, bold=True, color=TEXT_WHITE)
        # 설명
        add_text(slide, px + Inches(1.15), py + Inches(0.75), p_w - Inches(1.35), Inches(0.5),
                 p["desc"], font_size=11, color=TEXT_60)

# ════════════════════════════════════════════════════════════════════
# Slide 22 — Contact / 클로징
# ════════════════════════════════════════════════════════════════════

def slide_22_contact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_solid(slide, BG_PRIMARY)

    # 배경 그라데이션
    add_gradient_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT,
                      RGBColor(0x1e, 0x1b, 0x4b),
                      RGBColor(0x4a, 0x04, 0x4c),
                      angle=135)

    # 추상 원
    add_circle(slide, Inches(9.0), Inches(-1.5), Inches(7),
               fill_color=RGBColor(0x2e, 0x1a, 0x47))
    add_circle(slide, Inches(-3.0), Inches(4.5), Inches(6),
               fill_color=RGBColor(0x3a, 0x10, 0x40))

    # 중앙 상단 로고
    add_logo_placeholder(slide, Inches(5.0), Inches(0.7), Inches(3.3), Inches(0.9), size_pt=42)

    # 큰 카피
    add_text_multi(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.0), [
        ("마케터의 자연어 한 줄,", 46, True, TEXT_WHITE),
        ("AI가 마케팅 전체를 운영합니다.", 46, True, TEXT_WHITE),
    ], align=PP_ALIGN.CENTER, line_spacing=1.2)

    add_text(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.7),
             "지금 바로 시작하세요.",
             font_size=32, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)

    # CTA 버튼
    cta_y = Inches(4.95)
    cta_w = Inches(2.8)
    cta_h = Inches(0.7)
    cta_gap = Inches(0.3)
    cta_total = cta_w * 2 + cta_gap
    cta_start = (SLIDE_WIDTH - cta_total) / 2
    # primary
    add_rounded_gradient_rect(slide, cta_start, cta_y, cta_w, cta_h,
                              VIOLET, FUCHSIA, corner_radius_pct=0.40)
    add_text(slide, cta_start, cta_y, cta_w, cta_h,
             "무료 체험 시작 →",
             font_size=16, bold=True, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    # secondary
    add_rounded_rect(slide, cta_start + cta_w + cta_gap, cta_y, cta_w, cta_h,
                     fill_color=None, line_color=TEXT_WHITE, line_width=1.0,
                     corner_radius_pct=0.40)
    add_text(slide, cta_start + cta_w + cta_gap, cta_y, cta_w, cta_h,
             "영업 상담 신청",
             font_size=16, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Contact 영역
    add_text_multi(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9), [
        ("웹사이트  hanjul.ai", 14, False, TEXT_WHITE),
        ("고객사 관리자  app.hanjul.ai", 14, False, TEXT_WHITE),
        ("영업 문의  [이메일 / 전화]", 14, False, TEXT_60),
    ], align=PP_ALIGN.CENTER, line_spacing=1.5)

    # 하단 영문
    add_text(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
             "한줄로 — Where AI proposes, humans approve.",
             font_size=13, italic=True, color=VIOLET, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# 메인 흐름
# ════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_01_cover(prs)
    slide_02_manifesto(prs)
    slide_03_problem(prs)
    slide_04_solution(prs)
    slide_05_three_pillars(prs)
    slide_06_menu_1(prs)
    slide_07_menu_2(prs)
    slide_08_ai_flow(prs)
    slide_09_korea_native(prs)
    slide_10_cdp(prs)
    slide_11_cross_channel(prs)
    slide_12_mobile_dm(prs)
    slide_13_memory(prs)
    slide_14_compliance(prs)
    slide_15_cost_safety(prs)
    slide_16_payment(prs)
    slide_17_global_compare(prs)
    slide_18_pricing(prs)
    slide_19_results(prs)
    slide_20_onboarding(prs)
    slide_21_principles(prs)
    slide_22_contact(prs)

    output_path = "docs/한줄로_서비스소개서_v2_2026-05.pptx"
    prs.save(output_path)
    print(f"OK: {output_path}")
    print(f"슬라이드 수: {len(prs.slides)}")

if __name__ == "__main__":
    main()
